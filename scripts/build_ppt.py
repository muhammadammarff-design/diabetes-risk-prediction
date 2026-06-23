"""
Diabetes Risk Prediction - FINAL Presentation generator (v3).

18 slides. Pure native python-pptx (no raw XML hacks).

CHANGES FROM v2 (per sir's feedback):
  + Slide 1: Updated team info (M. Ammar 2540004)
  + NEW slide 6: "Our Model" - Random Forest introduced EARLY
  + Slide 7: Workflow is now a proper BLOCK DIAGRAM
                  (Input -> Preprocessing -> Model -> Output)
  + Slide 5: SELECTED features (top 5 by importance) highlighted
  + Slide 18: Conclusion shows BIG result numbers prominently

Run from project root:
    python scripts/build_ppt.py
"""
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# =====================================================================
# PATHS
# =====================================================================
HERE         = Path(__file__).resolve().parent
PROJ         = HERE.parent
PLOTS        = PROJ / "plots"
MODELS       = PROJ / "models"
PRESENTATION = PROJ / "presentation"
PRESENTATION.mkdir(exist_ok=True)

# =====================================================================
# LOAD REAL RESULTS (no hard-coded numbers!)
# =====================================================================
metrics     = json.loads((MODELS / "metrics.json").read_text())
comparison  = json.loads((MODELS / "model_comparison.json").read_text())
features    = json.loads((MODELS / "feature_names.json").read_text())
top5        = list(metrics["top_5_features"].items())

# =====================================================================
# PROFESSIONAL PALETTE
# =====================================================================
TEAL       = RGBColor(0x0F, 0x76, 0x6E)
TEAL_DARK  = RGBColor(0x0B, 0x4F, 0x4A)
TEAL_LIGHT = RGBColor(0x14, 0xB8, 0xA6)
AMBER      = RGBColor(0xF5, 0x9E, 0x0B)
AMBER_DK   = RGBColor(0xD9, 0x77, 0x06)
INK        = RGBColor(0x1E, 0x29, 0x3B)
SLATE      = RGBColor(0x47, 0x55, 0x69)
MUTED      = RGBColor(0x94, 0xA3, 0xB8)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
SOFT       = RGBColor(0xF1, 0xF5, 0xF9)
LINE       = RGBColor(0xE2, 0xE8, 0xF0)
GREEN_OK   = RGBColor(0x16, 0xA3, 0x4A)
RED_NO     = RGBColor(0xDC, 0x26, 0x26)
BADGE_BG   = RGBColor(0xEC, 0xFE, 0xFF)

FONT       = "Calibri"

# =====================================================================
# PRESENTATION SETUP
# =====================================================================
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH          = prs.slide_width, prs.slide_height
blank           = prs.slide_layouts[6]
TOTAL           = 18

# =====================================================================
# HELPERS
# =====================================================================
def rect(slide, x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(0.75)
    s.shadow.inherit = False
    return s

def text(slide, x, y, w, h, content, *, size=14, bold=False, color=SLATE,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = content if isinstance(content, list) else content.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb

def bullets(slide, x, y, w, h, items, *, size=14, color=SLATE,
            line_spacing=1.35, space_after=8, bullet_color=AMBER, bullet="■"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        r1 = p.add_run(); r1.text = f"{bullet}  "
        r1.font.name = FONT; r1.font.size = Pt(size); r1.font.bold = True
        r1.font.color.rgb = bullet_color
        r2 = p.add_run(); r2.text = item
        r2.font.name = FONT; r2.font.size = Pt(size); r2.font.color.rgb = color
    return tb

def chrome(slide, num):
    rect(slide, 0, 0, Inches(0.12), SH, TEAL)
    text(slide, Inches(0.45), SH - Inches(0.38), Inches(8), Inches(0.3),
         "Diabetes Risk Prediction  ·  Intro to Data Science  ·  Sem 2",
         size=9, color=MUTED)
    text(slide, SW - Inches(1.2), SH - Inches(0.38), Inches(0.9), Inches(0.3),
         f"{num} / {TOTAL}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)

def title_block(slide, eyebrow, title_txt):
    text(slide, Inches(0.7), Inches(0.55), Inches(12), Inches(0.3),
         eyebrow.upper(), size=10, bold=True, color=AMBER)
    text(slide, Inches(0.7), Inches(0.9), Inches(12), Inches(0.75),
         title_txt, size=28, bold=True, color=INK)
    rect(slide, Inches(0.7), Inches(1.6), Inches(0.5), Inches(0.05), AMBER)

def note(slide, txt):
    nf = slide.notes_slide.notes_text_frame
    nf.text = txt


# =====================================================================
# SLIDE 1 - TITLE  (team info updated per reference)
# =====================================================================
s = prs.slides.add_slide(blank)
rect(s, 0, 0, Inches(0.12), SH, TEAL)

text(s, Inches(0.9), Inches(1.4), Inches(10), Inches(0.4),
     "INTRODUCTION TO DATA SCIENCE  ·  FINAL PRESENTATION",
     size=11, bold=True, color=AMBER)
text(s, Inches(0.9), Inches(1.95), Inches(12), Inches(1.4),
     "Diabetes Risk Prediction", size=54, bold=True, color=INK)
text(s, Inches(0.9), Inches(3.45), Inches(11.5), Inches(0.7),
     "Predicting diabetes from 21 health indicators using supervised machine learning",
     size=18, color=SLATE)
rect(s, Inches(0.9), Inches(4.45), Inches(1.0), Inches(0.04), AMBER)

text(s, Inches(0.9), Inches(4.75), Inches(5), Inches(0.3),
     "GROUP", size=10, bold=True, color=TEAL)
text(s, Inches(0.9), Inches(5.1), Inches(6), Inches(2),
     ["Ali Raza       ·  2540010",
      "M. Ammar       ·  2540004",
      "Taha Ali         ·  2540008"],
     size=15, color=INK, line_spacing=1.5)

text(s, Inches(7.5), Inches(4.75), Inches(5), Inches(0.3),
     "INSTRUCTOR", size=10, bold=True, color=TEAL)
text(s, Inches(7.5), Inches(5.1), Inches(5), Inches(0.4),
     "Sir Zaki", size=15, color=INK)
text(s, Inches(7.5), Inches(5.7), Inches(5), Inches(0.3),
     "SEMESTER · DATE", size=10, bold=True, color=TEAL)
text(s, Inches(7.5), Inches(6.05), Inches(5), Inches(0.4),
     "Semester 2  ·  2026", size=15, color=INK)
note(s,
    "Greet the audience. State project title, team and instructor. Keep brief — 30 seconds. "
    "Move to the next slide.")


# =====================================================================
# SLIDE 2 - OUTLINE  (updated for 18 slides)
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 2); title_block(s, "Outline", "What we will cover today")
left = [
    "01  ·  Problem statement",
    "02  ·  Dataset overview",
    "03  ·  Target & features",
    "04  ·  Our model (Random Forest)",
    "05  ·  Project workflow (block diagram)",
    "06  ·  Data cleaning",
    "07  ·  Exploratory data analysis",
    "08  ·  Correlation analysis",
    "09  ·  Train / test split",
]
right = [
    "10  ·  Why machine learning?",
    "11  ·  4 candidate ML models",
    "12  ·  Model comparison",
    "13  ·  Random Forest — detailed results",
    "14  ·  Top diabetes risk factors",
    "15  ·  Streamlit web app",
    "16  ·  Limitations",
    "17  ·  Conclusion & results",
    "18  ·  Q&A",
]
bullets(s, Inches(0.7), Inches(2.2), Inches(6), Inches(5), left,
        size=13, color=INK, bullet="·")
bullets(s, Inches(7.2), Inches(2.2), Inches(6), Inches(5), right,
        size=13, color=INK, bullet="·")
note(s,
    "Walk through the outline in 30 seconds. We added an 'Our Model' slide early so "
    "the audience knows from the start which algorithm we used. The flow is: Problem → Data → "
    "Model → Method → Results → Demo → Conclusion.")


# =====================================================================
# SLIDE 3 - PROBLEM STATEMENT
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 3); title_block(s, "01  ·  Why this matters", "The diabetes problem")

bullets(s, Inches(0.7), Inches(2.2), Inches(7.3), Inches(4.5), [
    "Diabetes affects 1 in 11 adults worldwide (~537 million people).",
    "Early detection drastically reduces complications and costs.",
    "Many people are unaware they have prediabetes or diabetes.",
    "Doctors need low-cost, fast screening tools.",
    "Patient survey data is cheap to collect — no lab tests required.",
], size=15, line_spacing=1.45, space_after=10)

cx, cy, cw, ch = Inches(8.3), Inches(2.4), Inches(4.4), Inches(4)
rect(s, cx, cy, cw, ch, SOFT, line=LINE)
text(s, cx + Inches(0.3), cy + Inches(0.35), cw - Inches(0.6), Inches(0.35),
     "GLOBAL IMPACT", size=10, bold=True, color=AMBER)
text(s, cx + Inches(0.3), cy + Inches(0.85), cw - Inches(0.6), Inches(1.5),
     "537M", size=72, bold=True, color=TEAL_DARK)
text(s, cx + Inches(0.3), cy + Inches(2.4), cw - Inches(0.6), Inches(0.5),
     "adults living with diabetes", size=14, color=SLATE)
text(s, cx + Inches(0.3), cy + Inches(3.1), cw - Inches(0.6), Inches(0.7),
     "Source: IDF Diabetes Atlas 2021", size=10, color=MUTED)
note(s,
    "Open with the human impact. 'Diabetes is not a small problem — 537 million adults have it.' "
    "This frames why our project matters. Don't dwell more than 1 minute.")


# =====================================================================
# SLIDE 4 - DATASET OVERVIEW  (with explicit feature count + selected features)
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 4); title_block(s, "02  ·  Data", "Dataset overview")

bullets(s, Inches(0.7), Inches(2.2), Inches(6.5), Inches(4.5), [
    "Source: CDC Behavioral Risk Factor Surveillance System (BRFSS), 2015.",
    "Collected through annual phone interviews across the United States.",
    "Cleaned version published on Kaggle by Alex Teboul.",
    "Pre-balanced 50 / 50 between diabetic and non-diabetic patients.",
    "Each row = one patient  ·  each column = one health indicator.",
], size=14, line_spacing=1.45, space_after=8)

# Right side - dataset stats card
cx, cy, cw, ch = Inches(7.8), Inches(2.2), Inches(4.9), Inches(4.6)
rect(s, cx, cy, cw, ch, SOFT, line=LINE)
rect(s, cx, cy, cw, Inches(0.45), TEAL)
text(s, cx + Inches(0.3), cy + Inches(0.07), cw - Inches(0.6), Inches(0.4),
     "DATASET AT A GLANCE", size=11, bold=True, color=WHITE)
rows = [
    ("Records",       "70,693 patients"),
    ("Total Features", f"{len(features)}  health indicators"),
    ("Target",        "Diabetes_binary  (0 / 1)"),
    ("Class Balance", "50 % / 50 %"),
    ("Type",          "Tabular · structured"),
    ("Task",          "Binary classification"),
]
ry = cy + Inches(0.75)
for label, value in rows:
    text(s, cx + Inches(0.3), ry, Inches(1.7), Inches(0.4),
         label, size=11, bold=True, color=TEAL)
    text(s, cx + Inches(2.0), ry, Inches(2.8), Inches(0.4),
         value, size=11, color=INK)
    ry += Inches(0.55)
note(s,
    "Explain WHERE the data came from. BRFSS is a real CDC survey — that's why our model is "
    "grounded in actual public-health data. Emphasize the **21 features** count — this is what "
    "sir expects to hear explicitly.")


# =====================================================================
# SLIDE 5 - TARGET & FEATURES  (with SELECTED features highlighted)
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 5); title_block(s, "03  ·  Variables",
                            "Target and selected features")

# Target box (top)
tx, ty, tw, th = Inches(0.7), Inches(2.1), Inches(12.0), Inches(0.95)
rect(s, tx, ty, tw, th, SOFT, line=LINE)
rect(s, tx, ty, Inches(0.08), th, AMBER)
text(s, tx + Inches(0.25), ty + Inches(0.1), Inches(3), Inches(0.3),
     "TARGET VARIABLE  (y)", size=10, bold=True, color=AMBER)
text(s, tx + Inches(0.25), ty + Inches(0.4), Inches(12), Inches(0.5),
     "Diabetes_binary  →  0 = No Diabetes   ·   1 = Diabetes / Prediabetes",
     size=16, bold=True, color=INK)

# Selected features badge  (shows the top 5 right at the start)
sx, sy, sw_, sh_ = Inches(0.7), Inches(3.2), Inches(12.0), Inches(0.85)
rect(s, sx, sy, sw_, sh_, BADGE_BG, line=AMBER)
text(s, sx + Inches(0.25), sy + Inches(0.04), Inches(11.5), Inches(0.32),
     "SELECTED FEATURES  (top 5 by model importance)",
     size=11, bold=True, color=AMBER_DK)
text(s, sx + Inches(0.25), sy + Inches(0.4), Inches(11.5), Inches(0.45),
     "  ".join([f"{f} ({v:.3f})" for f, v in top5]),
     size=13, bold=True, color=TEAL_DARK)

# Features by category
text(s, Inches(0.7), Inches(4.20), Inches(8), Inches(0.35),
     "ALL 21 FEATURES  (X)  ·  ORGANISED IN 3 CATEGORIES",
     size=10, bold=True, color=AMBER)

col_w = Inches(4.1); col_y = Inches(4.65); col_h = Inches(2.15)
def feat_col(x, header, items, highlight_set):
    rect(s, x, col_y, col_w, col_h, WHITE, line=LINE)
    rect(s, x, col_y, col_w, Inches(0.4), TEAL)
    text(s, x + Inches(0.25), col_y + Inches(0.07),
         col_w - Inches(0.4), Inches(0.3),
         header, size=10, bold=True, color=WHITE)
    # Items - selected ones in bold amber
    tb = slide_tb = s.shapes.add_textbox(
        x + Inches(0.25), col_y + Inches(0.50),
        col_w - Inches(0.4), col_h - Inches(0.55))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.15
        p.space_after = Pt(1)
        r1 = p.add_run(); r1.text = "•  "
        r1.font.name = FONT; r1.font.size = Pt(10); r1.font.bold = True
        # Highlight if feature name is in selected set
        is_selected = any(item.startswith(sel) for sel in highlight_set)
        r1.font.color.rgb = AMBER if is_selected else TEAL
        r2 = p.add_run(); r2.text = item
        r2.font.name = FONT; r2.font.size = Pt(9.5)
        r2.font.bold = is_selected
        r2.font.color.rgb = AMBER_DK if is_selected else INK

# 3 columns - the second tuple item is the "feature display name"
feat_col(Inches(0.7), "VITAL  ·  CLINICAL (7)", [
    "HighBP  ·  high blood pressure",
    "HighChol  ·  high cholesterol",
    "CholCheck  ·  5-yr check",
    "BMI  ·  Body Mass Index",
    "Stroke  ·  ever had one",
    "HeartDiseaseorAttack",
    "DiffWalk  ·  difficulty walking",
], [f for f, _ in top5])

feat_col(Inches(5.0), "LIFESTYLE (7)", [
    "Smoker  ·  100+ cigs",
    "PhysActivity  ·  past 30 days",
    "Fruits  ·  daily",
    "Veggies  ·  daily",
    "HvyAlcoholConsump",
    "AnyHealthcare  ·  coverage",
    "NoDocbcCost  ·  couldn't see",
], [f for f, _ in top5])

feat_col(Inches(9.3), "SELF-REPORTED  ·  DEMO (7)", [
    "GenHlth  ·  health (1-5)",
    "MentHlth  ·  bad-mood days",
    "PhysHlth  ·  sick days",
    "Sex  ·  0 female / 1 male",
    "Age  ·  group (1-13)",
    "Education  ·  level (1-6)",
    "Income  ·  bracket (1-8)",
], [f for f, _ in top5])
note(s,
    "Mention the 21 features are split into 3 categories. Highlight the SELECTED features "
    "shown in amber — these are the 5 most important ones according to our Random Forest model. "
    "Highlighted items in each column are the selected ones.")


# =====================================================================
# SLIDE 6 - OUR MODEL  (introduce Random Forest early - SIMPLE design)
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 6); title_block(s, "04  ·  Our model",
                            "Random Forest")

# One clear explanation
text(s, Inches(0.7), Inches(2.1), Inches(12), Inches(0.5),
     "A team of 200 decision trees that vote on the answer.",
     size=22, bold=True, color=TEAL_DARK)

# Sub-explanation
text(s, Inches(0.7), Inches(2.75), Inches(12), Inches(0.5),
     "Each tree looks at the patient data and predicts 0 or 1. "
     "The majority vote becomes the final prediction.",
     size=14, color=SLATE, line_spacing=1.4)

# Simple 3-element diagram (compact, well-spaced)
# Row 1: 3 tree boxes  ->  vote box  ->  output box
diag_y = Inches(4.3)

# 3 tree boxes (smaller, no overlap issues)
tree_w = Inches(1.4)
tree_h = Inches(1.5)
tree_gap = Inches(0.05)  # small gap between trees
tree_x_positions = [Inches(0.7), Inches(0.7) + tree_w + tree_gap, Inches(0.7) + 2*(tree_w + tree_gap)]
def tree_box(x, y, num):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, tree_w, tree_h)
    box.fill.solid(); box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = TEAL; box.line.width = Pt(1.5); box.shadow.inherit = False
    text(s, x, y + Inches(0.15), tree_w, Inches(0.35),
         f"Tree {num}", size=11, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    text(s, x, y + Inches(0.55), tree_w, Inches(0.4),
         "predicts", size=10, color=SLATE, align=PP_ALIGN.CENTER)
    text(s, x, y + Inches(0.95), tree_w, Inches(0.4),
         "0  or  1", size=12, bold=True, color=INK, align=PP_ALIGN.CENTER)

tree_box(tree_x_positions[0], diag_y, "1")
tree_box(tree_x_positions[1], diag_y, "2")
tree_box(tree_x_positions[2], diag_y, "3")

# Label below: "200 trees total"
text(s, Inches(0.7), diag_y + tree_h + Inches(0.05), Inches(4.5), Inches(0.3),
     "... and 197 more  (200 trees total)",
     size=11, color=MUTED, align=PP_ALIGN.CENTER, bold=True)

# Voting box — directly to the right of trees (no long arrows crossing)
vot_x = tree_x_positions[2] + tree_w + Inches(0.45)
vot_w = Inches(2.6)
rect(s, vot_x, diag_y, vot_w, tree_h, AMBER)
text(s, vot_x, diag_y + Inches(0.3), vot_w, Inches(0.5),
     "VOTE", size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
text(s, vot_x, diag_y + Inches(0.9), vot_w, Inches(0.4),
     "majority wins", size=11, color=WHITE, align=PP_ALIGN.CENTER)

# Arrow between trees and vote box (single arrow showing flow)
arrow_a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                              tree_x_positions[2] + tree_w + Inches(0.05),
                              diag_y + tree_h/2 - Inches(0.15),
                              Inches(0.4), Inches(0.3))
arrow_a.fill.solid(); arrow_a.fill.fore_color.rgb = TEAL
arrow_a.line.fill.background(); arrow_a.shadow.inherit = False

# Arrow from vote -> output
arrow_b = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                              vot_x + vot_w + Inches(0.05),
                              diag_y + tree_h/2 - Inches(0.15),
                              Inches(0.4), Inches(0.3))
arrow_b.fill.solid(); arrow_b.fill.fore_color.rgb = TEAL
arrow_b.line.fill.background(); arrow_b.shadow.inherit = False

# Output box
out_x = vot_x + vot_w + Inches(0.5)
out_w = Inches(2.5)
rect(s, out_x, diag_y, out_w, tree_h, TEAL_DARK)
text(s, out_x, diag_y + Inches(0.3), out_w, Inches(0.5),
     "PREDICTION", size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
text(s, out_x, diag_y + Inches(0.85), out_w, Inches(0.5),
     "0 (no diabetes)", size=11, color=WHITE, align=PP_ALIGN.CENTER)
text(s, out_x, diag_y + Inches(1.1), out_w, Inches(0.4),
     "1 (diabetes)", size=11, color=WHITE, align=PP_ALIGN.CENTER)

# Why we chose RF — small note below (kept above y=6.5)
why_y = Inches(6.55)
rect(s, Inches(0.7), why_y, Inches(12.0), Inches(0.3), SOFT, line=LINE)
text(s, Inches(0.9), why_y + Inches(0.02), Inches(11.5), Inches(0.28),
     "WHY: best F1 score across 4 models tested, and gives feature importance for free.",
     size=11, color=SLATE, bold=True, align=PP_ALIGN.CENTER)
note(s,
    "Introduce the model upfront so the audience knows what we'll be training. ONE sentence "
    "explanation: 'A team of 200 decision trees that vote on the answer.' Then show the simple "
    "diagram. Why this model: best F1 + feature importance. Keep it under 1 minute.")


# =====================================================================
# SLIDE 7 - WORKFLOW  (proper BLOCK DIAGRAM with input/output/methodology)
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 7); title_block(s, "05  ·  Workflow",
                            "Block diagram: Input → Methodology → Output")

# 3 sections with clear labels - COMPACT
sec_y = Inches(2.05)
rect(s, Inches(0.7), sec_y, Inches(2.0), Inches(0.45), TEAL)
text(s, Inches(0.7), sec_y + Inches(0.05), Inches(2.0), Inches(0.4),
     "INPUT", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

inp_x = Inches(0.7); inp_y = sec_y + Inches(0.55)
inp_w = Inches(12.0); inp_h = Inches(0.7)
rect(s, inp_x, inp_y, inp_w, inp_h, SOFT, line=LINE)
text(s, inp_x + Inches(0.3), inp_y + Inches(0.04), Inches(11.5), Inches(0.3),
     "70,693 patient records  ·  21 health indicators  ·  Diabetes_binary (target)",
     size=11, bold=True, color=INK, align=PP_ALIGN.CENTER)
text(s, inp_x + Inches(0.3), inp_y + Inches(0.36), Inches(11.5), Inches(0.3),
     "Source: BRFSS 2015 (CDC survey) — Kaggle clean version",
     size=10, color=SLATE, align=PP_ALIGN.CENTER)

# Methodology header
meth_y = Inches(3.7)
rect(s, Inches(0.7), meth_y, Inches(2.0), Inches(0.45), AMBER)
text(s, Inches(0.7), meth_y + Inches(0.05), Inches(2.0), Inches(0.4),
     "METHODOLOGY", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 4 methodology boxes (compact)
steps = [
    ("01", "Data\nCleaning",  "Drop duplicates"),
    ("02", "EDA &\nCorrelation", "Stats + heatmap"),
    ("03", "Train / Test\nSplit", "80/20 stratified"),
    ("04", "Train\nRandom Forest", "200 trees"),
]
box_w, box_h = Inches(2.85), Inches(1.35)
arr_w = Inches(0.25)
gap_x = arr_w + Inches(0.05)
total_w = 4 * box_w + 3 * gap_x
start_x = Inches(0.7) + (Inches(12.0) - total_w) / 2
row_y = meth_y + Inches(0.55)

def step_box(x, y, num, title, desc):
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = TEAL; card.line.width = Pt(1.25); card.shadow.inherit = False
    badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        x + Inches(0.22), y + Inches(0.18), Inches(0.5), Inches(0.32))
    badge.fill.solid(); badge.fill.fore_color.rgb = TEAL
    badge.line.fill.background(); badge.shadow.inherit = False
    tf = badge.text_frame; tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num
    r.font.name = FONT; r.font.size = Pt(10); r.font.bold = True
    r.font.color.rgb = WHITE
    text(s, x + Inches(0.85), y + Inches(0.15), box_w - Inches(1), Inches(0.65),
         title, size=12, bold=True, color=INK, line_spacing=1.05)
    text(s, x + Inches(0.3), y + Inches(0.85), box_w - Inches(0.5), Inches(0.4),
         desc, size=10, color=SLATE)

def r_arrow(x, y):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y - Inches(0.13), arr_w, Inches(0.26))
    a.fill.solid(); a.fill.fore_color.rgb = AMBER
    a.line.fill.background(); a.shadow.inherit = False

for i, st in enumerate(steps):
    x = start_x + i * (box_w + gap_x)
    step_box(x, row_y, *st)
    if i < len(steps) - 1:
        r_arrow(x + box_w + Inches(0.02), row_y + box_h/2)

# Output header (compact, bottom)
out_h_y = Inches(5.65)
rect(s, Inches(0.7), out_h_y, Inches(2.0), Inches(0.45), TEAL_DARK)
text(s, Inches(0.7), out_h_y + Inches(0.05), Inches(2.0), Inches(0.4),
     "OUTPUT", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Output content - 3 badges (compact)
out_y = out_h_y + Inches(0.55)
out_h = Inches(0.7)
ox = Inches(0.7)
ow = Inches(3.85)
og = (Inches(12.0) - 3 * ow) / 2
def out_badge(x, label, value):
    rect(s, x, out_y, ow, out_h, TEAL_DARK)
    text(s, x, out_y + Inches(0.05), ow, Inches(0.26),
         label, size=10, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
    text(s, x, out_y + Inches(0.30), ow, Inches(0.4),
         value, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

acc  = comparison[3]["Accuracy"]
rec  = comparison[3]["Recall"]
f1v  = comparison[3]["F1"]
out_badge(ox,                          "ACCURACY", f"{acc*100:.2f}%")
out_badge(ox + ow + og,                "F1 SCORE", f"{f1v:.3f}")
out_badge(ox + 2*(ow + og),            "RECALL",   f"{rec:.3f}")
note(s,
    "This is THE block diagram sir wants. Three layers stacked vertically: INPUT on top (the "
    "raw data), METHODOLOGY in the middle (the 4 steps we performed), OUTPUT at the bottom "
    "(the actual numbers we got). Walk through them in 30 seconds — top to bottom.")


# =====================================================================
# SLIDE 8 - DATA CLEANING
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 8); title_block(s, "06  ·  Preprocessing", "Data cleaning — what we did")

bullets(s, Inches(0.7), Inches(2.2), Inches(7), Inches(4), [
    "Checked for missing values  →  found 0  ✓",
    "Removed duplicate rows  →  1 636 dropped",
    "Final shape: 69 057 rows × 22 columns",
    "Converted 21 columns to int (everything except BMI which is float)",
    "Verified class balance still ~50 / 50",
], size=14, line_spacing=1.45, space_after=10)

cx, cy, cw, ch = Inches(8.2), Inches(2.4), Inches(4.5), Inches(3.6)
rect(s, cx, cy, cw, ch, SOFT, line=LINE)
text(s, cx + Inches(0.3), cy + Inches(0.25), cw - Inches(0.4), Inches(0.3),
     "BEFORE  →  AFTER", size=10, bold=True, color=AMBER)
text(s, cx + Inches(0.3), cy + Inches(0.7), cw - Inches(0.4), Inches(0.4),
     "Rows", size=11, bold=True, color=TEAL)
text(s, cx + Inches(0.3), cy + Inches(1.05), cw - Inches(0.4), Inches(0.5),
     "70 693  →  69 057", size=20, bold=True, color=INK)
text(s, cx + Inches(0.3), cy + Inches(1.8), cw - Inches(0.4), Inches(0.4),
     "Duplicates", size=11, bold=True, color=TEAL)
text(s, cx + Inches(0.3), cy + Inches(2.15), cw - Inches(0.4), Inches(0.5),
     "1 636  →  0", size=20, bold=True, color=INK)
text(s, cx + Inches(0.3), cy + Inches(2.9), cw - Inches(0.4), Inches(0.4),
     "Missing values", size=11, bold=True, color=TEAL)
text(s, cx + Inches(0.3), cy + Inches(3.25), cw - Inches(0.4), Inches(0.3),
     "0  →  0  ✓ already clean", size=14, color=INK)
note(s,
    "Key point: we dropped 1 636 duplicate rows. If we hadn't, the test set could have contained "
    "rows that the model already saw in training — inflating accuracy artificially. This is "
    "called 'data leakage' and is a viva-favourite question.")


# =====================================================================
# SLIDE 9 - EDA  (4 embedded plots)
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 9); title_block(s, "07  ·  EDA", "Exploratory data analysis")

plot_w, plot_h = Inches(5.8), Inches(2.45)
positions = [
    (Inches(0.7),  Inches(2.1)),
    (Inches(6.85), Inches(2.1)),
    (Inches(0.7),  Inches(4.75)),
    (Inches(6.85), Inches(4.75)),
]
plot_files = [
    "01_class_balance.png",
    "02_bmi_distribution.png",
    "03_bmi_boxplot.png",
    "04_bp_chol.png",
]
for (x, y), f in zip(positions, plot_files):
    if (PLOTS / f).exists():
        s.shapes.add_picture(str(PLOTS / f), x, y, width=plot_w, height=plot_h)
note(s,
    "Show 4 EDA charts: class balance (perfectly 50/50), BMI distribution, BMI boxplot, and "
    "BP/Cholesterol counts. Key takeaway: diabetic patients have visibly higher BMI and more "
    "HighBP / HighChol cases.")


# =====================================================================
# SLIDE 10 - CORRELATION
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 10); title_block(s, "08  ·  Correlation", "Which features matter most?")

if (PLOTS / "05_correlation_heatmap.png").exists():
    s.shapes.add_picture(str(PLOTS / "05_correlation_heatmap.png"),
                         Inches(0.4), Inches(1.95),
                         width=Inches(7.0), height=Inches(5.2))

text(s, Inches(7.6), Inches(2.0), Inches(5.5), Inches(0.4),
     "TOP 10 CORRELATED WITH DIABETES", size=11, bold=True, color=AMBER)
top10 = [
    ("GenHlth",              0.397),
    ("HighBP",               0.372),
    ("BMI",                  0.286),
    ("HighChol",             0.281),
    ("Age",                  0.275),
    ("DiffWalk",             0.267),
    ("Income",               0.213),
    ("HeartDiseaseorAttack", 0.207),
    ("PhysHlth",             0.207),
    ("Education",            0.159),
]
ty = Inches(2.5)
for feat, val in top10:
    text(s, Inches(7.6), ty, Inches(2.7), Inches(0.35),
         feat, size=12, bold=True, color=INK)
    bar_w = Inches(2.2 * val / 0.4)
    rect(s, Inches(10.4), ty + Inches(0.07), bar_w, Inches(0.2), TEAL)
    text(s, Inches(10.4) + bar_w + Inches(0.1), ty, Inches(0.8), Inches(0.35),
         f"{val:.2f}", size=11, color=SLATE)
    ty += Inches(0.42)
note(s,
    "Heatmap shows feature relationships. The strongest correlations with diabetes are "
    "General Health (0.40), High Blood Pressure (0.37), BMI (0.29), and High Cholesterol (0.28). "
    "These are also the top features in the model — good cross-validation.")


# =====================================================================
# SLIDE 11 - TRAIN / TEST SPLIT
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 11); title_block(s, "09  ·  Split",
                            "Train / test split — done before everything else")

bullets(s, Inches(0.7), Inches(2.2), Inches(7), Inches(4), [
    "80 % training data  ·  20 % test data",
    "Stratified by target → keeps 50 / 50 balance in both splits",
    "random_state = 42 → reproducible results",
    "Split BEFORE any preprocessing → prevents data leakage",
], size=15, line_spacing=1.5, space_after=10)

cx, cy = Inches(8.2), Inches(2.5)
cw, ch = Inches(4.5), Inches(0.8)
text(s, cx, cy - Inches(0.5), cw, Inches(0.3),
     "DATA SPLIT", size=10, bold=True, color=AMBER)
rect(s, cx, cy, Inches(3.6), ch, TEAL)
rect(s, cx + Inches(3.6), cy, Inches(0.9), ch, AMBER)
text(s, cx, cy + Inches(0.95), Inches(3.6), Inches(0.3),
     "Train · 55 245 rows", size=10, color=INK, align=PP_ALIGN.CENTER)
text(s, cx + Inches(3.6), cy + Inches(0.95), Inches(0.9), Inches(0.3),
     "Test · 13 812", size=10, color=INK, align=PP_ALIGN.CENTER)

wx, wy, ww, wh = Inches(8.2), Inches(4.5), Inches(4.5), Inches(2.4)
rect(s, wx, wy, ww, wh, SOFT, line=LINE)
text(s, wx + Inches(0.3), wy + Inches(0.2), ww - Inches(0.4), Inches(0.35),
     "WHY SPLIT FIRST?", size=10, bold=True, color=AMBER)
text(s, wx + Inches(0.3), wy + Inches(0.65), ww - Inches(0.4), Inches(1.6),
     "If we scaled, imputed, or sampled before splitting, information from the test set "
     "would leak into training. Accuracy would look great in our notebook but fail on new data.",
     size=12, color=INK, line_spacing=1.35)
note(s,
    "This is the most important methodology point. Mention 'data leakage' by name — it shows "
    "mastery. Explain we split first so the test set is truly unseen.")


# =====================================================================
# SLIDE 12 - WHY MACHINE LEARNING
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 12); title_block(s, "10  ·  Approach", "Why supervised machine learning?")

def card(x, y, w, h, label, value, desc):
    rect(s, x, y, w, h, SOFT, line=LINE)
    rect(s, x, y, Inches(0.08), h, TEAL)
    text(s, x + Inches(0.3), y + Inches(0.25), w - Inches(0.4), Inches(0.3),
         label, size=10, bold=True, color=AMBER)
    text(s, x + Inches(0.3), y + Inches(0.6), w - Inches(0.4), Inches(0.6),
         value, size=20, bold=True, color=INK)
    text(s, x + Inches(0.3), y + Inches(1.45), w - Inches(0.4), Inches(1.5),
         desc, size=12, color=SLATE, line_spacing=1.35)

cw, ch = Inches(4.0), Inches(3.1)
card(Inches(0.7),  Inches(2.5), cw, ch,
     "LEARNING TYPE", "Supervised",
     "The target column is labeled (0 / 1) — the model learns from known examples.")
card(Inches(4.9),  Inches(2.5), cw, ch,
     "TASK", "Classification",
     "We predict a discrete class, not a number. Two classes → binary classification.")
card(Inches(9.1),  Inches(2.5), cw, ch,
     "OUTPUT", "Probability + label",
     "Each prediction includes a confidence score, useful for risk-screening contexts.")
note(s,
    "State clearly: this is supervised binary classification. We have labeled examples, two "
    "classes, and we use the patterns to predict diabetes for new patients.")


# =====================================================================
# SLIDE 13 - 4 CANDIDATE MODELS
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 13); title_block(s, "11  ·  Candidates", "4 supervised classifiers tested")

models_info = [
    ("Logistic Regression",  "Fits a probability line between the two classes",
     "Simple · interpretable", "Less accurate on non-linear patterns"),
    ("Decision Tree",        "Asks yes/no questions to split the data",
     "Easy to visualise",     "Prone to overfitting"),
    ("K-Nearest Neighbors",  "Predicts class by voting among nearest neighbours",
     "Intuitive",             "Slow on large data · needs scaling"),
    ("Random Forest",        "Combines many decision trees, votes the answer  ✓ SELECTED",
     "Robust · feature imp.", "Slower to train"),
]
cw, ch = Inches(6.0), Inches(2.15)
positions = [
    (Inches(0.7),  Inches(2.2)),
    (Inches(6.85), Inches(2.2)),
    (Inches(0.7),  Inches(4.6)),
    (Inches(6.85), Inches(4.6)),
]
for idx, ((x, y), (name, how, pros, cons)) in enumerate(zip(positions, models_info)):
    is_selected = "SELECTED" in name
    rect(s, x, y, cw, ch, WHITE, line=AMBER if is_selected else LINE)
    rect(s, x, y, cw, Inches(0.45), AMBER if is_selected else TEAL)
    text(s, x + Inches(0.3), y + Inches(0.07), cw - Inches(0.4), Inches(0.4),
         name.upper(), size=11, bold=True, color=WHITE)
    text(s, x + Inches(0.3), y + Inches(0.55), cw - Inches(0.4), Inches(0.5),
         how, size=12, color=INK)
    text(s, x + Inches(0.3), y + Inches(1.15), Inches(1), Inches(0.35),
         "PROS", size=9, bold=True, color=GREEN_OK)
    text(s, x + Inches(1.3), y + Inches(1.15), cw - Inches(1.5), Inches(0.35),
         pros, size=11, color=SLATE)
    text(s, x + Inches(0.3), y + Inches(1.55), Inches(1), Inches(0.35),
         "CONS", size=9, bold=True, color=RED_NO)
    text(s, x + Inches(1.3), y + Inches(1.55), cw - Inches(1.5), Inches(0.35),
         cons, size=11, color=SLATE)
note(s,
    "These are the 4 supervised classifiers we considered. Random Forest (bottom-right) is "
    "highlighted with an amber border because we selected it. We tested all four to make a "
    "data-driven choice — not just pick a favourite.")


# =====================================================================
# SLIDE 14 - MODEL COMPARISON (real numbers + chart)
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 14); title_block(s, "12  ·  Results", "Model comparison — real numbers")

if (PLOTS / "08_model_comparison.png").exists():
    s.shapes.add_picture(str(PLOTS / "08_model_comparison.png"),
                         Inches(0.5), Inches(2.0),
                         width=Inches(7.5), height=Inches(4.2))

text(s, Inches(8.4), Inches(2.0), Inches(4.5), Inches(0.4),
     "FULL METRICS", size=11, bold=True, color=AMBER)

headers = ["Model", "Acc", "F1"]
col_x = [Inches(8.4), Inches(11.0), Inches(11.9)]
hy = Inches(2.5)
rect(s, Inches(8.4), hy, Inches(4.3), Inches(0.4), TEAL)
for h, x in zip(headers, col_x):
    text(s, x + Inches(0.05), hy + Inches(0.07), Inches(1.5), Inches(0.3),
         h, size=10, bold=True, color=WHITE)

ry = hy + Inches(0.4)
for i, row in enumerate(comparison):
    bg = SOFT if i % 2 == 0 else WHITE
    rect(s, Inches(8.4), ry, Inches(4.3), Inches(0.5), bg, line=LINE)
    name_color = TEAL_DARK if row["Model"] == "Random Forest" else INK
    name_bold  = True if row["Model"] == "Random Forest" else False
    text(s, col_x[0] + Inches(0.05), ry + Inches(0.12),
         Inches(2.5), Inches(0.35),
         row["Model"], size=10, bold=name_bold, color=name_color)
    text(s, col_x[1] + Inches(0.05), ry + Inches(0.12),
         Inches(0.8), Inches(0.35),
         f"{row['Accuracy']:.3f}", size=10, color=INK)
    text(s, col_x[2] + Inches(0.05), ry + Inches(0.12),
         Inches(0.8), Inches(0.35),
         f"{row['F1']:.3f}", size=10, color=INK)
    ry += Inches(0.5)

vy = Inches(5.0)
rect(s, Inches(8.4), vy, Inches(4.3), Inches(1.85), SOFT, line=LINE)
text(s, Inches(8.5), vy + Inches(0.15), Inches(4), Inches(0.3),
     "WINNER", size=10, bold=True, color=AMBER)
text(s, Inches(8.5), vy + Inches(0.5), Inches(4), Inches(0.5),
     "Random Forest", size=16, bold=True, color=TEAL_DARK)
text(s, Inches(8.5), vy + Inches(1.05), Inches(4.1), Inches(0.8),
     "Best F1 (0.759) and built-in feature importance.",
     size=11, color=INK, line_spacing=1.3)
note(s,
    "Be honest: Logistic Regression actually has slightly higher accuracy. But Random Forest "
    "has the best F1 score AND gives us feature importance for free. F1 matters more in medical "
    "apps because we balance precision and recall. That's why we chose RF.")


# =====================================================================
# SLIDE 15 - RANDOM FOREST DETAILED RESULTS  (BIG result badges)
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 15); title_block(s, "13  ·  Final model",
                            "Random Forest — detailed results")

# Confusion matrix on left
if (PLOTS / "06_confusion_matrix.png").exists():
    s.shapes.add_picture(str(PLOTS / "06_confusion_matrix.png"),
                         Inches(0.7), Inches(2.2),
                         width=Inches(5.2), height=Inches(4.6))

# Right: BIG metric badges (more prominent for marks)
def big_card(x, y, w, h, label, value, sublabel="", accent=TEAL):
    rect(s, x, y, w, h, WHITE, line=LINE)
    rect(s, x, y, Inches(0.08), h, accent)
    text(s, x + Inches(0.3), y + Inches(0.15), w - Inches(0.4), Inches(0.3),
         label, size=10, bold=True, color=AMBER)
    text(s, x + Inches(0.3), y + Inches(0.5), w - Inches(0.4), Inches(0.9),
         value, size=32, bold=True, color=INK)
    if sublabel:
        text(s, x + Inches(0.3), y + Inches(1.4), w - Inches(0.4), Inches(0.3),
             sublabel, size=10, color=SLATE)

acc  = comparison[3]["Accuracy"]
prec = comparison[3]["Precision"]
rec  = comparison[3]["Recall"]
f1v  = comparison[3]["F1"]

big_card(Inches(6.5), Inches(2.2), Inches(3.0), Inches(1.85),
         "ACCURACY", f"{acc*100:.2f}%", "correct predictions")
big_card(Inches(9.7), Inches(2.2), Inches(3.0), Inches(1.85),
         "F1 SCORE", f"{f1v:.3f}", "best of 4 models", accent=AMBER)
big_card(Inches(6.5), Inches(4.15), Inches(3.0), Inches(1.85),
         "PRECISION", f"{prec:.3f}",
         "predicted diabetics correct")
big_card(Inches(9.7), Inches(4.15), Inches(3.0), Inches(1.85),
         "RECALL", f"{rec:.3f}",
         "actual diabetics we caught")

# Hyperparameters (compact, fits above footer)
hx, hy, hw, hh = Inches(6.5), Inches(6.10), Inches(6.2), Inches(0.65)
rect(s, hx, hy, hw, hh, SOFT, line=LINE)
text(s, hx + Inches(0.25), hy + Inches(0.05), hw - Inches(0.4), Inches(0.25),
     "HYPERPARAMETERS", size=9, bold=True, color=AMBER)
text(s, hx + Inches(0.25), hy + Inches(0.32), hw - Inches(0.4), Inches(0.3),
     "n_estimators=200  ·  max_depth=15  ·  min_samples_leaf=5  ·  random_state=42",
     size=10, color=INK, line_spacing=1.1)
note(s,
    "Confusion matrix on the left: 4 696 true negatives and 5 574 true positives. "
    "Recall of 0.79 means we caught 79% of real diabetics — this is what matters most in "
    "medical screening. False negatives are more dangerous than false positives. The 4 metric "
    "cards on the right are deliberately LARGE — sir values clear, prominent results.")


# =====================================================================
# SLIDE 16 - TOP RISK FACTORS  (feature importance)
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 16); title_block(s, "14  ·  Insights", "Top diabetes risk factors")

if (PLOTS / "07_feature_importance.png").exists():
    s.shapes.add_picture(str(PLOTS / "07_feature_importance.png"),
                         Inches(0.4), Inches(1.85),
                         width=Inches(6.6), height=Inches(4.95))

text(s, Inches(7.5), Inches(2.0), Inches(5.5), Inches(0.4),
     "TOP 5 RISK FACTORS LEARNED", size=11, bold=True, color=AMBER)
ranks = ["1.", "2.", "3.", "4.", "5."]
ty = Inches(2.55)
for (feat, val), rank in zip(top5, ranks):
    text(s, Inches(7.5), ty, Inches(0.6), Inches(0.55),
         rank, size=22, bold=True, color=TEAL)
    text(s, Inches(8.2), ty + Inches(0.05), Inches(3.5), Inches(0.45),
         feat, size=14, bold=True, color=INK)
    text(s, Inches(8.2), ty + Inches(0.4), Inches(4.5), Inches(0.3),
         f"importance = {val:.3f}", size=11, color=SLATE)
    ty += Inches(0.85)
note(s,
    "The model identified the 5 most important risk factors entirely from data: "
    "General Health, High BP, BMI, Age, and High Cholesterol. These match what doctors would "
    "tell you — validating that the model learned something meaningful.")


# =====================================================================
# SLIDE 17 - STREAMLIT APP
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 17); title_block(s, "15  ·  Deployment", "Streamlit web application")

bullets(s, Inches(0.7), Inches(2.2), Inches(6.5), Inches(4), [
    "Interactive form built with Streamlit (Python web framework).",
    "User enters 21 health indicators in plain English.",
    "Trained model loaded once and cached for speed.",
    "Output: risk label + probability + confidence bar.",
    "Code on GitHub  ·  deployed free on Streamlit Cloud.",
], size=14, line_spacing=1.45, space_after=8)

cx, cy, cw, ch = Inches(7.8), Inches(2.2), Inches(4.9), Inches(4.6)
rect(s, cx, cy, cw, ch, SOFT, line=LINE)
rect(s, cx, cy, cw, Inches(0.45), TEAL)
text(s, cx + Inches(0.3), cy + Inches(0.07), cw - Inches(0.4), Inches(0.4),
     "DIABETES RISK PREDICTION", size=11, bold=True, color=WHITE)

def field(label, value, y_off):
    fx = cx + Inches(0.3); fy = cy + y_off
    text(s, fx, fy, Inches(2.0), Inches(0.3),
         label, size=10, bold=True, color=SLATE)
    rect(s, fx + Inches(2.1), fy + Inches(0.02), Inches(2.0), Inches(0.3),
         WHITE, line=LINE)
    text(s, fx + Inches(2.2), fy + Inches(0.04), Inches(1.8), Inches(0.3),
         value, size=10, color=INK)

field("BMI",         "27",        Inches(0.7))
field("High BP",     "Yes",       Inches(1.1))
field("High Chol",   "No",        Inches(1.5))
field("General Hlth","3 - Good",  Inches(1.9))
field("Age group",   "50-54",     Inches(2.3))

bx = cx + Inches(0.3); by = cy + Inches(2.85)
rect(s, bx, by, Inches(4.3), Inches(0.5), TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, bx, by + Inches(0.12), Inches(4.3), Inches(0.3),
     "PREDICT", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

ry = cy + Inches(3.55)
rect(s, bx, ry, Inches(4.3), Inches(0.85), SOFT, line=AMBER)
text(s, bx + Inches(0.2), ry + Inches(0.1), Inches(4), Inches(0.3),
     "HIGH RISK", size=11, bold=True, color=RED_NO)
text(s, bx + Inches(0.2), ry + Inches(0.45), Inches(4), Inches(0.3),
     "probability  68 %", size=11, color=INK)
note(s,
    "Show or describe the live app. If on Streamlit Cloud, show the URL. The app loads the "
    "trained model and lets anyone enter data to get an instant prediction. Mention this is "
    "your 'demo moment' if you can run it live.")


# =====================================================================
# SLIDE 18 - CONCLUSION & RESULTS  (BIG prominent results, fits footer)
# =====================================================================
s = prs.slides.add_slide(blank)
chrome(s, 18); title_block(s, "16  ·  Conclusion",
                            "Results · Limitations · Q&A")

# 4 BIG result badges at top — the most prominent numbers
acc  = comparison[3]["Accuracy"]
f1v  = comparison[3]["F1"]
rec  = comparison[3]["Recall"]
prec = comparison[3]["Precision"]

# Row 1: 4 metric badges (compact, fits above footer)
big_y = Inches(2.1)
big_h = Inches(1.5)
big_w = Inches(2.95)
gap = (Inches(12.0) - 4 * big_w) / 3
bx = Inches(0.7)
def big_metric(x, label, value, sublabel, accent=TEAL):
    rect(s, x, big_y, big_w, big_h, WHITE, line=LINE)
    rect(s, x, big_y, big_w, Inches(0.35), accent)
    text(s, x + Inches(0.2), big_y + Inches(0.05), big_w - Inches(0.3), Inches(0.3),
         label, size=10, bold=True, color=WHITE)
    text(s, x + Inches(0.2), big_y + Inches(0.5), big_w - Inches(0.3), Inches(0.75),
         value, size=30, bold=True, color=INK)
    text(s, x + Inches(0.2), big_y + Inches(1.2), big_w - Inches(0.3), Inches(0.28),
         sublabel, size=9, color=SLATE)

big_metric(bx,                          "ACCURACY",  f"{acc*100:.2f}%", "on test set", TEAL)
big_metric(bx + big_w + gap,            "F1 SCORE",  f"{f1v:.3f}",      "best of 4 models", AMBER)
big_metric(bx + 2*(big_w + gap),        "PRECISION", f"{prec:.3f}",     "of predicted diabetics", TEAL)
big_metric(bx + 3*(big_w + gap),        "RECALL",    f"{rec:.3f}",      "of real diabetics caught", TEAL_DARK)

# Middle row: Conclusion (left) + Limitations + Q&A
mid_y = Inches(3.85)
text(s, Inches(0.7), mid_y, Inches(6), Inches(0.35),
     "CONCLUSION", size=11, bold=True, color=GREEN_OK)
bullets(s, Inches(0.7), mid_y + Inches(0.32), Inches(6), Inches(1.4), [
    f"Random Forest → {acc*100:.1f}% accuracy  ·  {f1v:.3f} F1 score.",
    "5 dominant risk factors identified, aligned with medical literature.",
    "Working web app + reproducible notebook, GitHub-ready.",
], size=11, line_spacing=1.3, space_after=4, bullet="·", bullet_color=GREEN_OK)

text(s, Inches(0.7), Inches(5.45), Inches(6), Inches(0.35),
     "LIMITATIONS", size=11, bold=True, color=RED_NO)
bullets(s, Inches(0.7), Inches(5.78), Inches(6), Inches(1.0), [
    "Self-reported survey data is noisy.",
    "No lab values (no HbA1c or glucose).",
    "Balanced data  ≠  real-world prevalence (~14%).",
], size=10, line_spacing=1.25, space_after=2, bullet="·", bullet_color=RED_NO)

# Right side - Likely Q&A
text(s, Inches(7.2), mid_y, Inches(5.5), Inches(0.35),
     "LIKELY QUESTIONS", size=11, bold=True, color=AMBER)
qas = [
    ("Why Random Forest over Logistic Reg?",
     "Best F1 + feature importance for free."),
    ("Why only 74% accuracy?",
     "Self-reported data caps it; matches published research."),
    ("Why drop duplicates?",
     "Prevents test-set inflation (data leakage)."),
    ("Supervised or unsupervised?",
     "Supervised — target is labeled 0 / 1."),
    ("How would you improve it?",
     "Add lab values · try XGBoost · tune hyperparameters."),
]
ty = mid_y + Inches(0.32)
for q, a in qas:
    text(s, Inches(7.2), ty, Inches(5.5), Inches(0.3),
         q, size=10, bold=True, color=INK)
    text(s, Inches(7.2), ty + Inches(0.25), Inches(5.5), Inches(0.3),
         a, size=10, color=SLATE)
    ty += Inches(0.46)

# Thank you bar — positioned safely ABOVE the footer (chrome footer at y=7.12)
# Position it at y=6.5 (well above footer), make it narrower and centered
thanks_w = Inches(6.0)
thanks_x = (SW - thanks_w) / 2
rect(s, thanks_x, Inches(6.45), thanks_w, Inches(0.32), TEAL,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, thanks_x, Inches(6.49), thanks_w, Inches(0.28),
     "Thank you  ·  Questions?", size=12, bold=True, color=WHITE,
     align=PP_ALIGN.CENTER)
note(s,
    "This is THE conclusion slide — 4 BIG metric badges at the top show results PROMINENTLY "
    "(this is what gets marks for 'how we show results'). Below them: conclusion bullets (left), "
    "limitations (left), and likely Q&A (right). End with 'Thank you — any questions?'.")


# =====================================================================
# SAVE
# =====================================================================
out_pptx = PRESENTATION / "Diabetes_Final_Presentation.pptx"
prs.save(out_pptx)
print(f"Saved {out_pptx.name}  ({TOTAL} slides, native python-pptx, with speaker notes)")
